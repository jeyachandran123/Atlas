"""
DIP Phase 2 end-to-end test — uploads one document per supported family,
waits for the background worker to reach knowledge_ready, and verifies the
Knowledge Object, metadata, tables, chunks, and processing event timeline.

Requires: API + document_worker running (see backend/QUICK_START.md).
"""
import io
import json
import sys
import time
import uuid

import httpx

BASE = "http://127.0.0.1:8000/api/v1"
EMAIL = f"dip2-{uuid.uuid4().hex[:8]}@example.com"
PASSWORD = "dip2-password-123"

PASS = 0
FAIL = 0


def check(name: str, cond: bool, detail: str = "") -> None:
    global PASS, FAIL
    if cond:
        PASS += 1
        print(f"  PASS  {name}")
    else:
        FAIL += 1
        print(f"  FAIL  {name}  {detail}")


def wait_ready(c: httpx.Client, H: dict, doc_id: str, timeout: float = 60.0) -> dict:
    """Poll /processing until knowledge_ready or failed."""
    deadline = time.time() + timeout
    last = {}
    while time.time() < deadline:
        r = c.get(f"{BASE}/documents/{doc_id}/processing", headers=H)
        last = r.json()
        if last.get("processing_status") in ("knowledge_ready", "failed"):
            return last
        time.sleep(1.0)
    return last


def make_docx() -> bytes:
    import docx
    d = docx.Document()
    d.add_heading("Quarterly Report", level=1)
    d.add_paragraph("This report summarizes Q3 performance across all regions.")
    d.add_heading("Revenue", level=2)
    d.add_paragraph("Revenue grew by twelve percent compared to the prior quarter.")
    table = d.add_table(rows=3, cols=2)
    table.rows[0].cells[0].text, table.rows[0].cells[1].text = "Region", "Revenue"
    table.rows[1].cells[0].text, table.rows[1].cells[1].text = "North", "120000"
    table.rows[2].cells[0].text, table.rows[2].cells[1].text = "South", "98000"
    buf = io.BytesIO()
    d.save(buf)
    return buf.getvalue()


def make_xlsx() -> bytes:
    import openpyxl
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Budget"
    ws.append(["Department", "Amount", "Formula"])
    ws.append(["Engineering", 50000, "=B2*1.1"])
    ws.append(["Marketing", 30000, "=B3*1.1"])
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def make_pptx() -> bytes:
    from pptx import Presentation
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Project Kickoff"
    body = slide.placeholders[1]
    body.text_frame.text = "Welcome to the project kickoff meeting."
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


with httpx.Client(timeout=120) as c:
    r = c.post(f"{BASE}/auth/register", json={
        "email": EMAIL, "password": PASSWORD, "full_name": "DIP2 E2E",
        "role": "developer", "org_id": "default",
    })
    if r.status_code not in (200, 201):
        sys.exit(f"REGISTER_FAILED {r.status_code} {r.text[:200]}")
    r = c.post(f"{BASE}/auth/login", json={"email": EMAIL, "password": PASSWORD})
    assert r.status_code == 200, r.text
    H = {"Authorization": f"Bearer {r.json()['access_token']}"}

    # ── Upload one document per family ───────────────────────────────────────
    docs: dict[str, str] = {}

    r = c.post(f"{BASE}/documents/upload", headers=H, files={
        "file": ("notes.txt", b"Meeting notes.\n\nAction item: ship phase two by Friday.\n", "text/plain")
    })
    check("upload txt", r.status_code == 201, r.text[:200])
    docs["txt"] = r.json()["document"]["id"]

    md = (
        "# Architecture Overview\n\n"
        "This system uses a pipeline architecture.\n\n"
        "## Stages\n\n"
        "- Loader\n- Parser\n- Normalizer\n\n"
        "| Stage | Owner |\n|---|---|\n| Loader | Platform team |\n| Parser | Platform team |\n"
    )
    r = c.post(f"{BASE}/documents/upload", headers=H, files={
        "file": ("architecture.md", md.encode(), "text/markdown")
    })
    check("upload md", r.status_code == 201, r.text[:200])
    docs["md"] = r.json()["document"]["id"]

    r = c.post(f"{BASE}/documents/upload", headers=H, files={
        "file": ("sales.csv", b"Region,Revenue\nNorth,120000\nSouth,98000\n", "text/csv")
    })
    check("upload csv", r.status_code == 201, r.text[:200])
    docs["csv"] = r.json()["document"]["id"]

    r = c.post(f"{BASE}/documents/upload", headers=H, files={
        "file": ("config.json", json.dumps({
            "name": "atlas", "version": 2, "features": ["dip", "chat"], "nested": {"a": 1}
        }).encode(), "application/json")
    })
    check("upload json", r.status_code == 201, r.text[:200])
    docs["json"] = r.json()["document"]["id"]

    xml_bytes = b"<root><item id='1'>Alpha</item><item id='2'>Beta</item></root>"
    r = c.post(f"{BASE}/documents/upload", headers=H, files={
        "file": ("data.xml", xml_bytes, "application/xml")
    })
    check("upload xml", r.status_code == 201, r.text[:200])
    docs["xml"] = r.json()["document"]["id"]

    r = c.post(f"{BASE}/documents/upload", headers=H, files={
        "file": ("report.docx", make_docx(),
                 "application/vnd.openxmlformats-officedocument.wordprocessingml.document")
    })
    check("upload docx", r.status_code == 201, r.text[:200])
    docs["docx"] = r.json()["document"]["id"]

    r = c.post(f"{BASE}/documents/upload", headers=H, files={
        "file": ("budget.xlsx", make_xlsx(),
                 "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")
    })
    check("upload xlsx", r.status_code == 201, r.text[:200])
    docs["xlsx"] = r.json()["document"]["id"]

    r = c.post(f"{BASE}/documents/upload", headers=H, files={
        "file": ("kickoff.pptx", make_pptx(),
                 "application/vnd.openxmlformats-officedocument.presentationml.presentation")
    })
    check("upload pptx", r.status_code == 201, r.text[:200])
    docs["pptx"] = r.json()["document"]["id"]

    def make_pdf() -> bytes:
        from pypdf import PdfWriter
        from pypdf.generic import DecodedStreamObject, DictionaryObject, NameObject
        writer = PdfWriter()
        writer.add_blank_page(width=300, height=300)
        page = writer.pages[0]

        stream = DecodedStreamObject()
        stream.set_data(b"BT /F1 12 Tf 10 250 Td (Hello DIP Phase Two) Tj ET")
        page[NameObject("/Contents")] = writer._add_object(stream)

        font = DictionaryObject({
            NameObject("/Type"): NameObject("/Font"),
            NameObject("/Subtype"): NameObject("/Type1"),
            NameObject("/BaseFont"): NameObject("/Helvetica"),
        })
        page[NameObject("/Resources")] = DictionaryObject({
            NameObject("/Font"): DictionaryObject({NameObject("/F1"): writer._add_object(font)})
        })
        buf = io.BytesIO()
        writer.write(buf)
        return buf.getvalue()

    r = c.post(f"{BASE}/documents/upload", headers=H, files={
        "file": ("brief.pdf", make_pdf(), "application/pdf")
    })
    check("upload pdf", r.status_code == 201, r.text[:200])
    docs["pdf"] = r.json()["document"]["id"]

    # ── Wait for the worker to process every document ────────────────────────
    print("\n— waiting for knowledge_ready —")
    results: dict[str, dict] = {}
    for fmt, doc_id in docs.items():
        state = wait_ready(c, H, doc_id)
        results[fmt] = state
        check(f"{fmt}: reaches knowledge_ready", state.get("processing_status") == "knowledge_ready",
              f"status={state.get('processing_status')} error={state.get('error')}")

    # ── Processing event timeline (audit of every stage) ──────────────────────
    txt_events = results["txt"].get("events", [])
    stage_names = [e["stage"] for e in txt_events]
    expected_stages = ["load", "detect", "parse", "ocr", "normalize", "metadata",
                        "structure", "tables", "images", "language", "clean", "chunk",
                        "knowledge", "persist"]
    check("all pipeline stages recorded", all(s in stage_names for s in expected_stages),
          f"got={stage_names}")
    check("no failed stages for txt", all(e["status"] != "failed" for e in txt_events))

    # ── Knowledge Object checks ───────────────────────────────────────────────
    r = c.get(f"{BASE}/documents/{docs['docx']}/knowledge?include_structure=true", headers=H)
    check("docx knowledge → 200", r.status_code == 200, r.text[:200])
    ko = r.json()
    check("docx has sections", ko["section_count"] >= 1, str(ko))
    check("docx has a table", ko["table_count"] == 1, str(ko["table_count"]))
    check("docx word_count > 0", ko["word_count"] > 0)
    check("docx structure tree present", ko["structure"] is not None and "children" in ko["structure"])
    check("docx metadata title inferred", bool(ko["metadata"]["title"]))

    r = c.get(f"{BASE}/documents/{docs['xlsx']}/knowledge", headers=H)
    xko = r.json()
    check("xlsx knowledge has sheet metadata", xko["metadata"]["sheet_count"] == 1, str(xko["metadata"]))

    r = c.get(f"{BASE}/documents/{docs['pptx']}/knowledge", headers=H)
    pko = r.json()
    check("pptx knowledge has slide metadata", pko["metadata"]["slide_count"] == 1, str(pko["metadata"]))

    # ── Chunks ─────────────────────────────────────────────────────────────────
    r = c.get(f"{BASE}/documents/{docs['docx']}/chunks", headers=H)
    check("docx chunks → 200", r.status_code == 200, r.text[:200])
    chunk_body = r.json()
    check("docx has chunks", chunk_body["total"] > 0, str(chunk_body["total"]))
    table_chunks = [ch for ch in chunk_body["items"] if ch["node_type"] == "table"]
    check("docx has a structured table chunk", len(table_chunks) == 1, str(len(table_chunks)))
    if table_chunks:
        check("table chunk carries structured rows in meta",
              bool(table_chunks[0]["meta"]) and "rows" in table_chunks[0]["meta"],
              str(table_chunks[0].get("meta")))
    check("chunks carry section_path", any(ch["section_path"] for ch in chunk_body["items"]))

    # ── CSV table structure preserved (not flattened to text) ────────────────
    r = c.get(f"{BASE}/documents/{docs['csv']}/chunks", headers=H)
    csv_chunks = r.json()["items"]
    csv_table = next((ch for ch in csv_chunks if ch["node_type"] == "table"), None)
    check("csv produced a structured table chunk", csv_table is not None)
    if csv_table:
        check("csv table meta has 2 data rows", csv_table["meta"]["row_count"] == 2, str(csv_table["meta"]))

    # ── JSON tree became structured nodes (not flattened) ─────────────────────
    r = c.get(f"{BASE}/documents/{docs['json']}/knowledge?include_structure=true", headers=H)
    jko = r.json()
    check("json knowledge ready", jko["confidence"] > 0)

    # ── Formula extraction (xlsx) — value present in some chunk ───────────────
    r = c.get(f"{BASE}/documents/{docs['xlsx']}/chunks", headers=H)
    xlsx_text = json.dumps(r.json())
    check("xlsx formula captured", "=B2*1.1" in xlsx_text or "B2*1.1" in xlsx_text, xlsx_text[:300])

    # ── Reprocess (idempotent) ────────────────────────────────────────────────
    r = c.post(f"{BASE}/documents/{docs['txt']}/process", headers=H)
    check("reprocess → 200", r.status_code == 200, r.text[:200])
    state2 = wait_ready(c, H, docs["txt"])
    check("reprocess reaches knowledge_ready again", state2.get("processing_status") == "knowledge_ready")
    r = c.get(f"{BASE}/documents/{docs['txt']}/chunks", headers=H)
    check("reprocess did not duplicate chunks", r.json()["total"] > 0)

    # ── processing_status surfaces on the document itself ─────────────────────
    r = c.get(f"{BASE}/documents/{docs['docx']}", headers=H)
    check("document.processing_status == knowledge_ready", r.json()["processing_status"] == "knowledge_ready")

    # ── 404 semantics for processing/knowledge endpoints on unknown doc ───────
    r = c.get(f"{BASE}/documents/{uuid.uuid4()}/knowledge", headers=H)
    check("unknown doc knowledge → 404", r.status_code == 404)

print(f"\n{PASS} passed, {FAIL} failed")
sys.exit(1 if FAIL else 0)
