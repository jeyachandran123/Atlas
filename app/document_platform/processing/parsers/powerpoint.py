"""PowerPoint (.pptx) parser — slides, titles, text, images, speaker notes."""
from __future__ import annotations

import io

from loguru import logger

from app.document_platform.processing.capabilities import ParserCapabilities
from app.document_platform.processing.models import (
    DocumentNode, ImageRef, NodeType, ParsedDocument, RawMetadata,
)
from app.document_platform.processing.parsers.base import AbstractDocumentParser, ParserError


class PowerPointParser(AbstractDocumentParser):
    name = "powerpoint"
    extensions = (".pptx",)
    version = "1.0.0"
    capabilities = ParserCapabilities(
        supports_tables=False,
        supports_images=True,
        supports_structure=True,  # title + slide sequence
    )

    def parse(self, content: bytes, filename: str) -> ParsedDocument:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        try:
            prs = Presentation(io.BytesIO(content))
        except Exception as e:
            raise ParserError(f"Cannot open .pptx: {e}") from e

        root = DocumentNode(type=NodeType.DOCUMENT)
        images: list[ImageRef] = []

        for slide_no, slide in enumerate(prs.slides, start=1):
            slide_node = root.add(DocumentNode(type=NodeType.SLIDE, page=slide_no))

            title_shape = getattr(slide.shapes, "title", None)
            if title_shape is not None and (title_shape.text or "").strip():
                slide_node.add(DocumentNode(
                    type=NodeType.HEADING, text=title_shape.text.strip(),
                    level=1, page=slide_no,
                ))

            for shape in slide.shapes:
                if shape is title_shape:
                    continue
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        img = shape.image
                        images.append(ImageRef(
                            name=img.filename or f"slide{slide_no}-img",
                            content=img.blob,
                            page=slide_no,
                            format=(img.ext or "png").lower(),
                        ))
                    except Exception as e:
                        logger.debug(f"pptx image extraction failed on slide {slide_no}: {e}")
                    continue
                if shape.has_text_frame:
                    text = "\n".join(
                        p.text for p in shape.text_frame.paragraphs if p.text.strip()
                    ).strip()
                    if text:
                        slide_node.add(DocumentNode(
                            type=NodeType.PARAGRAPH, text=text, page=slide_no,
                        ))

            if slide.has_notes_slide:
                notes = (slide.notes_slide.notes_text_frame.text or "").strip()
                if notes:
                    slide_node.add(DocumentNode(
                        type=NodeType.NOTE, text=notes, page=slide_no,
                        meta={"kind": "speaker_notes"},
                    ))

        props = prs.core_properties
        return ParsedDocument(
            root=root,
            images=images,
            raw_metadata=RawMetadata(
                title=props.title or "",
                author=props.author or "",
                created=str(props.created or ""),
                modified=str(props.modified or ""),
                slide_count=len(prs.slides),
            ),
            parser_name=self.name,
        )
